#!/usr/bin/env python3
"""Generate presentation/output.pptx from presentation/slides-gstack.json."""

from __future__ import annotations

import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError:
    print("Missing dependency: pip install -r requirements-pptx.txt", file=sys.stderr)
    sys.exit(1)

# RestoDesk theme — matches app/DESIGN.md
BG = RGBColor(0x1A, 0x14, 0x10)
CARD = RGBColor(0x2A, 0x22, 0x1C)
ACCENT = RGBColor(0xE0, 0x7A, 0x3A)
HIGHLIGHT = RGBColor(0xF5, 0xC5, 0x42)
TEXT = RGBColor(0xF5, 0xF0, 0xEB)
MUTED = RGBColor(0xA8, 0x9F, 0x94)

SLIDE_W = Inches(10)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.5)
TITLE_BAR_H = Inches(0.85)


def find_project_root(start: Path) -> Path:
    candidates = [start, *start.parents]
    for p in candidates:
        if (p / "presentation").is_dir():
            return p
    return start


def load_slides(project_root: Path) -> dict:
    pres = project_root / "presentation"
    candidates = [
        pres / "slides.json",
        pres / "slides-gstack.json",
        pres / "template-slides-gstack-5min.json",
        pres / "template-slides.json",
    ]
    for path in candidates:
        if path.exists():
            if path.name.startswith("template"):
                print(f"Using template: {path}")
            return json.loads(path.read_text(encoding="utf-8"))
    print(
        f"Error: no slides JSON in {pres}. "
        "Create slides-gstack.json from template-slides-gstack-5min.json.",
        file=sys.stderr,
    )
    sys.exit(1)


def set_slide_bg(slide, color: RGBColor = BG) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def resolve_image(project_root: Path, rel_path: str | None) -> Path | None:
    if not rel_path:
        return None
    img_path = project_root / rel_path
    return img_path if img_path.exists() else None


def add_speaker_notes(slide, notes: str | None) -> None:
    if not notes:
        return
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.text = notes.strip()


def add_title_bar(slide, title: str) -> None:
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, TITLE_BAR_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = CARD
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(MARGIN, Inches(0.12), Inches(9), Inches(0.65))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(26)
    tp.font.bold = True
    tp.font.color.rgb = ACCENT


def add_bullets(
    slide,
    bullets: list[str],
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int = 16,
) -> None:
    if not bullets:
        return
    body = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = body.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT
        p.space_after = Pt(8)


def add_picture_fit(slide, img_path: Path, left, top, width, height=None) -> None:
    if height is not None:
        slide.shapes.add_picture(str(img_path), left, top, width=width, height=height)
    else:
        slide.shapes.add_picture(str(img_path), left, top, width=width)


def add_opening_title_slide(
    prs: Presentation,
    title: str,
    subtitle: str,
    project_root: Path,
    hero_image: str | None = None,
    speaker_notes: str | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    img_path = resolve_image(project_root, hero_image)
    if img_path:
        add_picture_fit(slide, img_path, MARGIN, Inches(4.2), Inches(9), height=Inches(2.8))

    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(1.4))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER

    sub = slide.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(8.4), Inches(0.9))
    sp = sub.text_frame.paragraphs[0]
    sp.text = subtitle
    sp.font.size = Pt(22)
    sp.font.color.rgb = MUTED
    sp.alignment = PP_ALIGN.CENTER

    add_speaker_notes(slide, speaker_notes)


def layout_image_hero(slide, slide_data: dict, project_root: Path) -> None:
    add_title_bar(slide, slide_data.get("title", "Slide"))

    img_path = resolve_image(project_root, slide_data.get("image"))
    if img_path:
        add_picture_fit(slide, img_path, MARGIN, Inches(1.05), Inches(9), height=Inches(4.0))

    bullets = slide_data.get("bullets", [])
    add_bullets(slide, bullets, 0.5, 5.2, 9.0, 1.8, font_size=15)


def layout_split(slide, slide_data: dict, project_root: Path) -> None:
    add_title_bar(slide, slide_data.get("title", "Slide"))

    bullets = slide_data.get("bullets", [])
    add_bullets(slide, bullets, 0.5, 1.1, 4.3, 5.8, font_size=15)

    images = slide_data.get("images") or slide_data.get("screenshots") or []
    if slide_data.get("image") and not images:
        images = [slide_data["image"]]

    y_positions = [1.1, 3.85]
    img_h = Inches(2.5)
    img_w = Inches(4.5)
    for i, rel_path in enumerate(images[:2]):
        img_path = resolve_image(project_root, rel_path)
        if img_path:
            add_picture_fit(slide, img_path, Inches(5.0), Inches(y_positions[i]), img_w, height=img_h)


def layout_image_bottom(slide, slide_data: dict, project_root: Path) -> None:
    add_title_bar(slide, slide_data.get("title", "Slide"))

    bullets = slide_data.get("bullets", [])
    add_bullets(slide, bullets, 0.5, 1.05, 9.0, 0.9, font_size=16)

    img_path = resolve_image(project_root, slide_data.get("image"))
    if img_path:
        add_picture_fit(slide, img_path, MARGIN, Inches(2.1), Inches(9), height=Inches(5.0))


def layout_default(slide, slide_data: dict, project_root: Path) -> None:
    """Fallback for legacy bullet + screenshots format."""
    add_title_bar(slide, slide_data.get("title", "Slide"))
    bullets = slide_data.get("bullets", [])
    add_bullets(slide, bullets, 0.5, 1.1, 9.0, 3.5, font_size=16)

    screenshots = slide_data.get("screenshots", [])
    if screenshots:
        y = Inches(4.8)
        max_w = Inches(2.8)
        for i, rel_path in enumerate(screenshots[:3]):
            img_path = resolve_image(project_root, rel_path)
            if img_path:
                slide.shapes.add_picture(str(img_path), Inches(0.8 + i * 3.0), y, width=max_w)


LAYOUT_HANDLERS = {
    "image_hero": layout_image_hero,
    "split": layout_split,
    "image_bottom": layout_image_bottom,
    "default": layout_default,
}


def add_content_slide(prs: Presentation, slide_data: dict, project_root: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    layout = slide_data.get("layout", "default")
    handler = LAYOUT_HANDLERS.get(layout, layout_default)
    handler(slide, slide_data, project_root)

    add_speaker_notes(slide, slide_data.get("speaker_notes"))


def build_pptx(project_root: Path, data: dict) -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    add_opening_title_slide(
        prs,
        data.get("title", "AI ADAPT BATTLE"),
        data.get("subtitle", "Demo Presentation"),
        project_root,
        hero_image=data.get("hero_image"),
        speaker_notes=data.get("title_speaker_notes"),
    )

    for slide_data in data.get("slides", []):
        add_content_slide(prs, slide_data, project_root)

    out_dir = project_root / "presentation"
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / "output.pptx"
    prs.save(str(out_path))
    return out_path


def main() -> None:
    script_dir = Path(__file__).resolve().parent
    if script_dir.name == "scripts" and (script_dir.parent / "presentation").is_dir():
        project_root = script_dir.parent
    else:
        project_root = find_project_root(Path.cwd())

    if len(sys.argv) > 1:
        project_root = Path(sys.argv[1]).resolve()

    data = load_slides(project_root)
    out = build_pptx(project_root, data)
    slide_count = 1 + len(data.get("slides", []))
    print(f"Created {out} ({slide_count} slides)")


if __name__ == "__main__":
    main()
